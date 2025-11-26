"""
File processing service
Handles file uploads, validation, and text extraction
"""

import os
import tempfile
from typing import Tuple
from fastapi import UploadFile, HTTPException

# Text extraction libraries
import pdfplumber
from docx import Document
from pptx import Presentation
import openpyxl


class FileService:
    ALLOWED_TYPES = {'.pptx', '.docx', '.pdf', '.xlsx'}
    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB

    @staticmethod
    def validate_file(file: UploadFile) -> Tuple[str, int]:
        """
        Validate uploaded file
        Returns: (file_extension, file_size)
        """
        # Check filename
        if not file.filename:
            raise HTTPException(400, "No filename provided")

        # Check extension
        file_ext = os.path.splitext(file.filename)[1].lower()
        if file_ext not in FileService.ALLOWED_TYPES:
            raise HTTPException(
                400,
                f"File type {file_ext} not allowed. Allowed: {', '.join(FileService.ALLOWED_TYPES)}"
            )

        # Check file size (if available)
        if hasattr(file, 'size') and file.size:
            if file.size > FileService.MAX_FILE_SIZE:
                raise HTTPException(400, f"File size exceeds 10MB limit")

        return file_ext, file.size if hasattr(file, 'size') else 0

    @staticmethod
    async def extract_text(file_path: str, file_type: str) -> str:
        """
        Extract text from uploaded file
        """
        try:
            if file_type == '.pdf':
                return FileService._extract_pdf(file_path)
            elif file_type == '.docx':
                return FileService._extract_docx(file_path)
            elif file_type == '.pptx':
                return FileService._extract_pptx(file_path)
            elif file_type == '.xlsx':
                return FileService._extract_xlsx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            raise HTTPException(500, f"Text extraction failed: {str(e)}")

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        """Extract text from PDF"""
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return '\n'.join(text)

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        """Extract text from DOCX"""
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs if para.text])

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return '\n'.join(text)

    @staticmethod
    def _extract_xlsx(file_path: str) -> str:
        """Extract text from XLSX"""
        wb = openpyxl.load_workbook(file_path)
        text = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell])
                if row_text:
                    text.append(row_text)
        return '\n'.join(text)

    @staticmethod
    async def save_temp_file(file: UploadFile) -> str:
        """Save uploaded file to temporary location"""
        temp_file = tempfile.NamedTemporaryFile(delete=False)
        try:
            content = await file.read()
            temp_file.write(content)
            temp_file.close()
            return temp_file.name
        except Exception as e:
            os.unlink(temp_file.name)
            raise HTTPException(500, f"Failed to save file: {str(e)}")
